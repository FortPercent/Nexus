"""上传文件预处理：xlsx / csv / zip → 转 markdown 或递归解压后再送 Letta。
其他格式 (pdf / docx / txt / md) 原样透传由 Letta 自己处理。
"""
import csv as _csv
import datetime as _dt
import io
import logging
import os
import zipfile
from typing import List, Tuple

from fastapi import HTTPException


ProcessedFile = Tuple[str, bytes, str]  # (filename, content_bytes, mime_type)

ZIP_MAX_DEPTH = 3
ZIP_MAX_FILES = 50
ZIP_MAX_UNCOMPRESSED = 200 * 1024 * 1024  # 200 MB

PASSTHROUGH_EXTS = {"pdf", "txt", "md", "png", "jpg", "jpeg"}  # Letta 原生接受的 blob

_MIME_BY_EXT = {
    "pdf": "application/pdf",
    "txt": "text/plain",
    "md": "text/markdown",
    "png": "image/png",
    "jpg": "image/jpeg",
    "jpeg": "image/jpeg",
}


def _ext(filename: str) -> str:
    return (filename.rsplit(".", 1)[-1] or "").lower() if "." in filename else ""


def _fmt_cell(v) -> str:
    """markdown 表格单元格格式化：日期 / 浮点 / 管道符转义"""
    if v is None:
        return ""
    if isinstance(v, _dt.datetime):
        # 没时间部分就只留日期
        if v.hour == 0 and v.minute == 0 and v.second == 0:
            return v.strftime("%Y-%m-%d")
        return v.strftime("%Y-%m-%d %H:%M:%S")
    if isinstance(v, _dt.date):
        return v.strftime("%Y-%m-%d")
    if isinstance(v, float):
        # 1.0 → "1"，1.5 → "1.5"
        if v.is_integer():
            return str(int(v))
        return f"{v:g}"
    s = str(v)
    # markdown 表格里 | 会破表；\n 换成空格避免换行断表
    return s.replace("|", "\\|").replace("\n", " ").replace("\r", "")


MAX_ROWS_PER_SHEET = 5000  # 单 sheet 行上限；50MB 上传限额已是硬防线，这里放宽避免业务表被砍半


def _xlsx_to_markdown(data: bytes) -> str:
    """Streaming 版: 不 list(ws.iter_rows(...)), 大表 RSS 保持常数级.
    之前 list() 一次性吃进内存, 50MB × 并发时 RSS 尖峰. openpyxl iter_rows
    本身就是 generator, 只要不 list() + 达到 MAX_ROWS_PER_SHEET 立刻 break."""
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    out = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        out.append(f"## 📑 工作表: {sheet}")
        row_iter = ws.iter_rows(values_only=True)

        # 找第一个非空行当 header
        header = None
        for row in row_iter:
            if any(c is not None for c in row):
                header = row
                break
        if header is None:
            out.append("_(空表)_")
            out.append("")
            continue

        headers = [_fmt_cell(c) for c in header]
        out.append("| " + " | ".join(headers) + " |")
        out.append("|" + "---|" * len(headers))

        # 流式读 body, 满 MAX_ROWS 就 break
        written = 0
        extra_seen = 0  # 超限后只计数, 不写 out
        for row in row_iter:
            if written < MAX_ROWS_PER_SHEET:
                cells = [_fmt_cell(c) for c in row]
                while len(cells) < len(headers):
                    cells.append("")
                cells = cells[: len(headers)]
                out.append("| " + " | ".join(cells) + " |")
                written += 1
            else:
                extra_seen += 1
                # 再扫一小段确认真有额外行就停 — 避免无谓遍历剩余百万行
                if extra_seen > 100:
                    out.append(f"_(…另有 {MAX_ROWS_PER_SHEET}+ 行已省略, 前 {MAX_ROWS_PER_SHEET} 行即为样本)_")
                    break
        else:
            if extra_seen > 0:
                out.append(f"_(…另有 {extra_seen} 行已省略)_")
        out.append("")
    wb.close()
    return "\n".join(out)


def _docx_to_markdown(data: bytes) -> str:
    """python-docx 按节读 docx：标题/正文/表格转 markdown。图片/批注/tracked changes 不处理。"""
    try:
        import docx as _docx
    except ImportError:
        raise HTTPException(500, "服务端未装 python-docx")
    doc = _docx.Document(io.BytesIO(data))
    out = []

    # 遍历 paragraphs + tables（按文档顺序）
    # docx body 下元素包含 w:p 和 w:tbl，用 body._element 原生顺序
    body = doc.element.body
    tbl_iter = iter(doc.tables)

    def _para_to_md(para):
        text = para.text.strip()
        if not text:
            return None
        s = (para.style.name or "").lower()
        if "heading 1" in s or s == "title":
            return f"# {text}"
        if "heading 2" in s:
            return f"## {text}"
        if "heading 3" in s:
            return f"### {text}"
        if "heading 4" in s:
            return f"#### {text}"
        if "list" in s or "bullet" in s:
            return f"- {text}"
        return text

    def _table_to_md(table):
        if not table.rows:
            return None
        header = [_fmt_cell(c.text) for c in table.rows[0].cells]
        lines = ["", "| " + " | ".join(header) + " |", "|" + "---|" * len(header)]
        for row in table.rows[1:]:
            cells = [_fmt_cell(c.text) for c in row.cells]
            while len(cells) < len(header):
                cells.append("")
            lines.append("| " + " | ".join(cells[: len(header)]) + " |")
        lines.append("")
        return "\n".join(lines)

    # 按 body 顺序遍历
    para_idx = 0
    tbl_idx = 0
    for child in body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "p":
            if para_idx < len(doc.paragraphs):
                md = _para_to_md(doc.paragraphs[para_idx])
                if md:
                    out.append(md)
                para_idx += 1
        elif tag == "tbl":
            if tbl_idx < len(doc.tables):
                md = _table_to_md(doc.tables[tbl_idx])
                if md:
                    out.append(md)
                tbl_idx += 1

    if not out:
        out.append("_(文档无可提取内容；可能仅含图片或表单)_")
    return "\n".join(out)


def _pptx_to_markdown(data: bytes) -> str:
    """pptx 用 python-pptx 抽文本: 每张 slide 的标题 + 文本 + 表格 + notes. 图片/SmartArt 忽略."""
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(500, "服务端未装 python-pptx")
    prs = Presentation(io.BytesIO(data))
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"## 📑 Slide {idx}")
        # 先出标题 (如果 layout 有 title placeholder)
        title_text = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
        except Exception:
            pass
        if title_text:
            out.append(f"### {title_text}")
        # 所有非 title 的文本框 + 表格
        for shape in slide.shapes:
            # skip title shape 已处理
            if shape == getattr(slide.shapes, "title", None):
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = (para.text or "").strip()
                    if text:
                        out.append(text)
            elif shape.has_table:
                tbl = shape.table
                rows = list(tbl.rows)
                if not rows:
                    continue
                header = [_fmt_cell(c.text) for c in rows[0].cells]
                out.append("")
                out.append("| " + " | ".join(header) + " |")
                out.append("|" + "---|" * len(header))
                for row in rows[1:]:
                    cells = [_fmt_cell(c.text) for c in row.cells]
                    while len(cells) < len(header):
                        cells.append("")
                    out.append("| " + " | ".join(cells[: len(header)]) + " |")
                out.append("")
        # 演讲者备注
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    out.append(f"_演讲者备注: {notes}_")
        except Exception:
            pass
        out.append("")
    if not out or all(not line.strip() for line in out):
        out.append("_(无可提取文本; 可能全是图片/SmartArt)_")
    return "\n".join(out)


# 全局 semaphore: 防并发 4 worker 同时跑 libreoffice 爆内存 (150MB × 4 = 600MB)
# threading.Semaphore 是同步的, 用 threading.BoundedSemaphore 更稳 (release 过多会报错)
import threading as _threading
_LIBREOFFICE_SEM = _threading.BoundedSemaphore(2)  # 同一个 adapter 容器最多 2 并发


def _cleanup_stale_lo_tempdirs():
    """adapter 启动时清理 /tmp 下残留的 lo_* 目录 (上次转换崩溃 / kill -9 留的).
    和 kill 任何残留 soffice 进程."""
    import glob, shutil, subprocess as _sp
    for d in glob.glob("/tmp/lo_*"):
        try:
            shutil.rmtree(d, ignore_errors=True)
        except Exception:
            pass
    try:
        _sp.run(["pkill", "-9", "-f", "soffice"], capture_output=True, timeout=5)
    except Exception:
        pass


def _convert_via_libreoffice(data: bytes, src_ext: str, to_ext: str) -> bytes:
    """用 libreoffice headless 转老 Office 格式 → 新格式 bytes.
    保护层:
      1. BoundedSemaphore 限并发 2 (防 4 worker 全跑 libreoffice 打爆内存)
      2. TemporaryDirectory 自动清临时文件
      3. 60s 硬超时 + TimeoutExpired 时 kill 子进程组 (防 zombie)
      4. 非 0 返回码 → 明确 400 而非 500
    """
    import subprocess
    import tempfile
    import signal

    if not _LIBREOFFICE_SEM.acquire(timeout=30):
        # 30s 排队超时: 当前并发满且排队超限, 明确拒绝不要堆积
        raise HTTPException(503, "libreoffice 转换队列繁忙，请稍后重试")
    try:
        with tempfile.TemporaryDirectory(prefix="lo_") as tmpdir:
            in_path = os.path.join(tmpdir, f"input.{src_ext}")
            with open(in_path, "wb") as f:
                f.write(data)
            profile = f"file://{tmpdir}/profile"
            cmd = [
                "libreoffice",
                f"-env:UserInstallation={profile}",
                "--headless", "--nologo", "--nofirststartwizard",
                "--convert-to", to_ext,
                "--outdir", tmpdir,
                in_path,
            ]
            # start_new_session=True 让子进程独立进程组, timeout kill 时可整组 killpg
            proc = subprocess.Popen(
                cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                start_new_session=True,
            )
            try:
                stdout, stderr = proc.communicate(timeout=60)
            except subprocess.TimeoutExpired:
                # 整个 soffice 进程组一起 kill, 防 zombie
                try:
                    os.killpg(proc.pid, signal.SIGKILL)
                except ProcessLookupError:
                    pass
                proc.wait(timeout=5)
                logging.warning(f"libreoffice convert timeout src={src_ext}")
                raise HTTPException(400, f"libreoffice 转换超时 (60s), 文件可能过大或损坏")
            if proc.returncode != 0:
                logging.warning(f"libreoffice convert failed rc={proc.returncode} stderr={stderr[:300]!r}")
                raise HTTPException(400, f"libreoffice 转换失败 (rc={proc.returncode})")
            out_path = os.path.join(tmpdir, f"input.{to_ext}")
            if not os.path.exists(out_path):
                raise HTTPException(400, "libreoffice 未生成输出文件")
            with open(out_path, "rb") as f:
                return f.read()
    finally:
        _LIBREOFFICE_SEM.release()


def _csv_to_markdown(data: bytes) -> str:
    """Streaming 版: 逐行读 + 到 MAX_ROWS break. 大 CSV 不 list()."""
    # 试多种编码
    text = None
    for enc in ("utf-8", "gbk", "utf-16"):
        try:
            text = data.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = data.decode("utf-8", errors="replace")
    reader = _csv.reader(io.StringIO(text))

    # 读 header
    try:
        header = next(reader)
    except StopIteration:
        return "_(空 CSV)_"

    headers = [_fmt_cell(c) for c in header]
    out = ["| " + " | ".join(headers) + " |", "|" + "---|" * len(headers)]

    written = 0
    extra_seen = 0
    for row in reader:
        if written < MAX_ROWS_PER_SHEET:
            cells = [_fmt_cell(c) for c in row]
            while len(cells) < len(headers):
                cells.append("")
            out.append("| " + " | ".join(cells[: len(headers)]) + " |")
            written += 1
        else:
            extra_seen += 1
            if extra_seen > 100:
                out.append(f"_(…另有 {MAX_ROWS_PER_SHEET}+ 行已省略, 前 {MAX_ROWS_PER_SHEET} 行即为样本)_")
                break
    else:
        if extra_seen > 0:
            out.append(f"_(…另有 {extra_seen} 行已省略)_")
    return "\n".join(out)


def _unzip(data: bytes, parent_name: str, depth: int, counter: List[int]) -> List[ProcessedFile]:
    if depth > ZIP_MAX_DEPTH:
        raise HTTPException(400, f"zip 嵌套超过 {ZIP_MAX_DEPTH} 层")
    results: List[ProcessedFile] = []
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            total = sum(zi.file_size for zi in z.infolist())
            if total > ZIP_MAX_UNCOMPRESSED:
                raise HTTPException(413, f"zip 解压后 {total // 1024 // 1024}MB 超过 {ZIP_MAX_UNCOMPRESSED // 1024 // 1024}MB 限制")
            for zi in z.infolist():
                if zi.is_dir():
                    continue
                # 跳过隐藏 / __MACOSX
                if zi.filename.startswith("__MACOSX/") or os.path.basename(zi.filename).startswith("."):
                    continue
                counter[0] += 1
                if counter[0] > ZIP_MAX_FILES:
                    raise HTTPException(400, f"zip 内文件数超过 {ZIP_MAX_FILES}")
                with z.open(zi) as f:
                    inner = f.read()
                sub_name = f"{parent_name.rsplit('.zip', 1)[0]}__{zi.filename}"
                # 递归处理（可能是嵌套 zip）
                results.extend(_process(sub_name, inner, depth=depth + 1, counter=counter))
    except zipfile.BadZipFile:
        raise HTTPException(400, "zip 文件损坏或非 zip 格式")
    return results


def _process(filename: str, data: bytes, depth: int = 0, counter=None) -> List[ProcessedFile]:
    counter = counter if counter is not None else [0]
    ext = _ext(filename)

    if ext in PASSTHROUGH_EXTS:
        return [(filename, data, _MIME_BY_EXT.get(ext, "application/octet-stream"))]

    if ext == "xls":
        logging.info(f"LEGACY_OFFICE_CONVERT ext=.xls filename={filename}")
        try:
            xlsx_bytes = _convert_via_libreoffice(data, "xls", "xlsx")
            md = _xlsx_to_markdown(xlsx_bytes)
            return [(filename + ".md", md.encode("utf-8"), "text/markdown")]
        except HTTPException as e:
            logging.warning(f".xls convert fail: {filename}: {e.detail}")
            raise

    if ext == "xlsx":
        try:
            md = _xlsx_to_markdown(data)
        except Exception as e:
            logging.warning(f"xlsx 提取失败 {filename}: {e}")
            raise HTTPException(400, f"xlsx 解析失败: {e}")
        return [(filename + ".md", md.encode("utf-8"), "text/markdown")]

    if ext == "csv":
        try:
            md = _csv_to_markdown(data)
        except Exception as e:
            raise HTTPException(400, f"csv 解析失败: {e}")
        return [(filename + ".md", md.encode("utf-8"), "text/markdown")]

    if ext == "docx":
        try:
            md = _docx_to_markdown(data)
        except HTTPException:
            raise
        except Exception as e:
            logging.warning(f"docx 提取失败 {filename}: {e}")
            raise HTTPException(400, f"docx 解析失败: {e}")
        return [(filename + ".md", md.encode("utf-8"), "text/markdown")]

    if ext == "doc":
        logging.info(f"LEGACY_OFFICE_CONVERT ext=.doc filename={filename}")
        try:
            docx_bytes = _convert_via_libreoffice(data, "doc", "docx")
            md = _docx_to_markdown(docx_bytes)
            return [(filename + ".md", md.encode("utf-8"), "text/markdown")]
        except HTTPException as e:
            logging.warning(f".doc convert fail: {filename}: {e.detail}")
            raise

    if ext == "ppt":
        logging.info(f"LEGACY_OFFICE_CONVERT ext=.ppt filename={filename}")
        try:
            pptx_bytes = _convert_via_libreoffice(data, "ppt", "pptx")
            md = _pptx_to_markdown(pptx_bytes)
            return [(filename + ".md", md.encode("utf-8"), "text/markdown")]
        except HTTPException as e:
            logging.warning(f".ppt convert fail: {filename}: {e.detail}")
            raise

    if ext == "pptx":
        try:
            md = _pptx_to_markdown(data)
        except HTTPException:
            raise
        except Exception as e:
            logging.warning(f"pptx 提取失败 {filename}: {e}")
            raise HTTPException(400, f"pptx 解析失败: {e}")
        return [(filename + ".md", md.encode("utf-8"), "text/markdown")]

    if ext == "zip":
        return _unzip(data, filename, depth=depth, counter=counter)

    raise HTTPException(400, f"不支持的文件类型: .{ext}（支持 pdf/docx/txt/md/xlsx/csv/pptx/ppt/png/jpg/jpeg/zip）")


def process_upload(filename: str, data: bytes) -> List[ProcessedFile]:
    """入口：返回要上传到 Letta 的文件列表（pdf/docx/md 原样透传；xlsx/csv 转 md；zip 展开）"""
    if not data:
        raise HTTPException(400, "文件为空")
    return _process(filename, data)
